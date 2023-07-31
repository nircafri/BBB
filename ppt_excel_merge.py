import glob
import pandas as pd
import numpy as np
import glob
import datetime
from pptx import Presentation
from pptx.util import Inches

# load Losartan_PSWE_analysis_eeg_mri_time.csv
Losartan_PSWE_analysis_eeg_mri_time = pd.read_csv('Losartan_PSWE_analysis_eeg_mri_time.csv')
# add column for Time_in_PSWE_std
Losartan_PSWE_analysis_eeg_mri_time['Time_in_PSWE_Zscore'] = (Losartan_PSWE_analysis_eeg_mri_time['Time_in_PSWE'] - Losartan_PSWE_analysis_eeg_mri_time['Time_in_PSWE'].mean())/Losartan_PSWE_analysis_eeg_mri_time['Time_in_PSWE'].std(ddof=0)
# load Epilepsy_individulas.pptx
prs = Presentation('Epilepsy_individulas.pptx')
# load Epilepsy losartan eeg and mri dates.csv
Epilepsy_clinical_data = pd.read_excel('Epilepsy_clinical_data.xlsx', sheet_name='dates')
bool_first_slide = True
initials = ''
# run over each slide
for slide in prs.slides:
    if initials != slide.shapes[0].text:
        bool_first_slide = True
    initials = slide.shapes[0].text
    try:
        Session = int(initials.split('_')[1])
    except:
        continue
    columns = ['Time_in_PSWE','Mean_duration','Recording_length','Delta_MRI_EEG_Days']
    Losartan_PSWE_analysis_eeg_mri_time_rows = Losartan_PSWE_analysis_eeg_mri_time[Losartan_PSWE_analysis_eeg_mri_time['Session'] == Session]
    if Losartan_PSWE_analysis_eeg_mri_time_rows.shape[0] == 0:
        continue
    # get row in Epilepsy_clinical_data of Session in column code
    date_mri = Epilepsy_clinical_data[Epilepsy_clinical_data['Initials'].str.contains(str(Session))]['Scan Date'].values[0]
    # remove ' from date_mri
    date_mri = date_mri.replace("'", "")
    # to datetime
    date_mri = datetime.datetime.strptime(date_mri, '%Y.%m.%d')
    #  Losartan_PSWE_analysis_eeg_mri_time_rows['Recording_date'] to datetime
    Losartan_PSWE_analysis_eeg_mri_time_rows['Recording_date'] = Losartan_PSWE_analysis_eeg_mri_time_rows['Recording_date'].apply(lambda x: datetime.datetime.strptime(x, '%d.%m.%y'))
    # new column Delta_MRI_EEG_Days absolute value of date_mri - Recording_date
    Losartan_PSWE_analysis_eeg_mri_time_rows['Delta_MRI_EEG_Days'] = Losartan_PSWE_analysis_eeg_mri_time_rows['Recording_date'].apply(lambda x: abs(date_mri - x))
    # Losartan_PSWE_analysis_eeg_mri_time_rows = min(Losartan_PSWE_analysis_eeg_mri_time_rows['Delta_MRI_EEG_Days'])
    Losartan_PSWE_analysis_eeg_mri_time_rows = Losartan_PSWE_analysis_eeg_mri_time_rows[Losartan_PSWE_analysis_eeg_mri_time_rows['Delta_MRI_EEG_Days'] == min(Losartan_PSWE_analysis_eeg_mri_time_rows['Delta_MRI_EEG_Days'])].iloc[0].T
    # Losartan_PSWE_analysis_eeg_mri_time_rows[['Time_in_PSWE_zscore','Time_in_PSWE','Mean_duration','Median_duration'Fp1_time_in_events	Fp2_time_in_events	F3_time_in_events	F4_time_in_events	C3_time_in_events	C4_time_in_events	P3_time_in_events	P4_time_in_events	O1_time_in_events	O2_time_in_events	F7_time_in_events	F8_time_in_events	T3_time_in_events	T4_time_in_events	T5_time_in_events	T6_time_in_events	Fz_time_in_events	Cz_time_in_events	Pz_time_in_events
    Losartan_PSWE_analysis_eeg_mri_time_rows = Losartan_PSWE_analysis_eeg_mri_time_rows[columns] #,'Fp1_time_in_events','Fp2_time_in_events','F3_time_in_events','F4_time_in_events','C3_time_in_events','C4_time_in_events','P3_time_in_events','P4_time_in_events','O1_time_in_events','O2_time_in_events','F7_time_in_events','F8_time_in_events','T3_time_in_events','T4_time_in_events','T5_time_in_events','T6_time_in_events','Fz_time_in_events','Cz_time_in_events','Pz_time_in_events']]
    df = pd.DataFrame(columns=columns)
    df.loc[0] = Losartan_PSWE_analysis_eeg_mri_time_rows
    if bool_first_slide:
        bool_first_slide = False
        # add table Losartan_PSWE_analysis_eeg_mri_time_rows to slide
        shapes = slide.shapes
        x, y, cx, cy = Inches(1), Inches(3), Inches(10), Inches(1)
        table_placeholder = shapes.add_table(df.shape[0]+1, df.shape[1], x, y, cx, cy).table
        table = table_placeholder
        # write header row
        for i in range(df.shape[1]):
            table.cell(0,i).text = df.columns[i]
        # write data rows
        for i in range(df.shape[0]):
            for j in range(df.shape[1]):
                table.cell(i+1,j).text = str(df.values[i,j])
# save pptx
prs.save('Epilepsy_individulas2.pptx')
        


